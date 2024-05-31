# System/App Dependencies
import streamlit as st
import sys
import subprocess
import os
from dotenv import load_dotenv

# Files Readers
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation


from langchain.text_splitter import CharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.llms import HuggingFaceHub,LlamaCpp

from langchain_community.embeddings import HuggingFaceInstructEmbeddings
from langchain.memory import ConversationBufferMemory  # LLM Memory
from langchain.chains import ConversationalRetrievalChain   # LLM Chain
import tempfile
from langchain_community.document_loaders import CSVLoader

load_dotenv()
api_token = os.getenv("HUGGINGFACEHUB_API_TOKEN")



def Get_files_Text(files):
    text = ""
    for file in files: 
        file_type = file.name.split(".")[-1]

        if file_type == "pdf":
            pdf_reader = PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text()

        elif file_type == "docx":
            doc = Document(file)
            for para in doc.paragraphs:
                text += para.text

        elif file_type == "pptx":
            ppt = Presentation(file)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

        elif file_type == "txt":
            text += file.read().decode('utf-8') + "\n"

    return text

def get_text_chunks(raw_text):
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1000,
        chunk_overlap=100,
        length_function=len
    )
    text_chunks = text_splitter.split_text(raw_text)
    return text_chunks




def get_Vector_Store(text_chunks):
    embeddings = HuggingFaceInstructEmbeddings(model_name="hkunlp/instructor-xl")
    vector_store = FAISS.from_texts(texts=text_chunks,embedding=embeddings)
    return vector_store


def get_conversation_chain(vector_store,llm_model):

    llm = HuggingFaceHub(repo_id=llm_model,
                          model_kwargs={"temperature" : 0.2,"max_length":512},
                          task="text2text-generation",
                            huggingfacehub_api_token=api_token) 


# if you want to run locally using Mistral model uncomment the following code: -  
#     llm = LlamaCpp(
#     model_path=r"Models\mistral-7b-instruct-v0.2.Q4_K_S.gguf",
#     temperature=0.5,
    # n_ctx=1048 
#     # top_p=1, 
#     # verbose=True,
#     
# ) 
            

    memory = ConversationBufferMemory(
    memory_key="chat_history",
    output_key="answer",
      return_messages=True
    ) 
    
    conversation_chain = ConversationalRetrievalChain.from_llm(llm=llm, chain_type='stuff',
                                                 retriever=vector_store.as_retriever(search_kwargs={"k": 2}),
                                                 memory=memory)
 
    return conversation_chain



# >> for mistral model
def handle_user_input(user_query):
    response = st.session_state.Conversation.invoke({"question": user_query, "chat_history": st.session_state.chat_history})
    
    st.session_state.chat_history = response["chat_history"]
    return response["answer"]

def get_csv_chain(db,llm_model):  # No memory for csv files (gives better responses without memory)
    llm = HuggingFaceHub(repo_id=llm_model,
                                        model_kwargs={"temperature" : 0.1},
                                        huggingfacehub_api_token=api_token)
    
    
    
    chain = ConversationalRetrievalChain.from_llm(llm=llm, chain_type='stuff',
                            retriever=db.as_retriever())
    return chain




def main():
    
    st.set_page_config(page_title="Chatgbt Ghost", page_icon=":ghost:")
    st.title("Chatgbt Ghost")
    # Initialization: -
    if "Conversation" not in st.session_state:
        st.session_state.Conversation = None
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []
    if "llm_model" not in st.session_state:
        llm_model ="None"
        st.session_state.llm_model = None
    if  "vector_store" not in st.session_state:
        st.session_state.vector_store = None



    with st.sidebar:
        
        st.header("Upload files 📁")
        Uploaded_Files=st.file_uploader(" ", accept_multiple_files=True,type=['pdf', 'txt', 'docx','csv',"pptx"])

    if Uploaded_Files:
        option = st.radio("Use Ghost for: -",["Text-based Files","CSV Files"],key="Upload_Files")
    file_col, Model_col = st.columns(2)

    

    with Model_col:
        llm_model = st.selectbox(
            "Choose a Model",
            ("None","google/flan-t5-xxl","google/flan-t5-small", "google/flan-t5-base"),index=0, key="model")
        
        # Change model in mid conversation!
        if st.session_state.vector_store:
                    if st.session_state.llm_model != llm_model and llm_model != "None":
                        if st.session_state.llm_model != None:
                            if st.button("Reset and Start Model!",key="Create_Conversation_Chain") and st.session_state.vector_store:
                                st.session_state.llm_model = llm_model
                                st.session_state.chat_history = []  
                                if option == "CSV Files":
                                    st.session_state.Conversation = get_csv_chain(st.session_state.vector_store, st.session_state.llm_model)
                                else:  
                                    st.session_state.Conversation = get_conversation_chain(st.session_state.vector_store,st.session_state.llm_model)
                                
                        else:
                            st.session_state.llm_model = llm_model
                            if option == "CSV Files":
                                st.session_state.Conversation = get_csv_chain(st.session_state.vector_store, st.session_state.llm_model)
                            else:
                                st.session_state.Conversation = get_conversation_chain(st.session_state.vector_store,st.session_state.llm_model)
                                
        
        

    with file_col:
        if Uploaded_Files:
            files_dict = {f.name: f for f in Uploaded_Files}
            non_csv_filenames = [f.name for f in Uploaded_Files if not f.name.endswith('.csv')]
            csv_filenames = [f.name for f in Uploaded_Files if f.name.endswith('.csv')]
            if option == "Text-based Files" and non_csv_filenames:
                selected_files = st.multiselect(f"Select from {option}: -", non_csv_filenames)
                # get file object from its name
                selected_files = [files_dict[f] for f in selected_files]
                if selected_files:
                    if st.button("Process Files",key="Process_Files"):
                        raw_text = Get_files_Text(selected_files) 
                        # Get text Chunks from pdf
                        text_chunks = get_text_chunks(raw_text)
                        #st.write(text_chunks)
                        with st.status("Embedding Information ...") as s:
                            st.session_state.vector_store = get_Vector_Store(text_chunks)
                            s.update(label="Successfully Embedded Document Information!")

            elif option == "CSV Files" and csv_filenames:
                # CSV MODE: -
                
                # choose only 1 csv file
                selected_file = st.selectbox(f"Select from {option}: -", csv_filenames)
                selected_file = files_dict[selected_file]
                if selected_file:
                    if st.button("Process CSV File", key="Process_CSV_Files"):
                        if llm_model == "None":
                            st.error("Please Choose a Model First!")
                        else:
                            with st.status("Embedding Information ...") as s:
                                with tempfile.NamedTemporaryFile(delete=False) as tmp_file:
                                    tmp_file.write(selected_file.getvalue())
                                    tmp_file_path = tmp_file.name
                                loader = CSVLoader(file_path=tmp_file_path,encoding="utf-8",csv_args={"delimiter":","})
                                data = loader.load()
                                embeddings = HuggingFaceInstructEmbeddings(model_name="hkunlp/instructor-xl")
                                st.session_state.vector_store = FAISS.from_documents(data,embeddings)
                                s.update(label="Successfully Embedded CSV Information!")
            else: 
                st.write(f"There are no {option} files uploaded")
                                
                                         

        else:   
            st.write("No Uploaded Files Yet!")

    
            
# show the csv name with a heading
    if st.session_state.llm_model:
        st.success(f"{st.session_state.llm_model} Model is Ready! 🚀")    

    st.write("--------------------------")
    if st.session_state.chat_history:   
        for i, message in enumerate(st.session_state.chat_history):
            if i % 2 == 0:
                with st.chat_message("Human"):
                    st.markdown(message.content)
            else:
                with st.chat_message("ai"):
                    st.markdown(message.content)  

    user_query = st.chat_input("Type here...")
    if user_query is not None and user_query !="": 
        with st.chat_message("human"):
            st.markdown(user_query)

        bot_response=handle_user_input(user_query)

        with st.chat_message("ai"):
            st.markdown(bot_response)


# main() # for streamlit run Chatgbt_Ghost.py



if __name__ == '__main__':
    if len(sys.argv) > 1 and sys.argv[1] == 'run':
        main()
    else:
        subprocess.run(['streamlit', 'run', sys.argv[0], 'run'])
